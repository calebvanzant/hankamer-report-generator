"""
generate_report.py
------------------
Reads a filled-in data file (XLSX or CSV) and produces a finished
Hankamer info-session PowerPoint in the output/ folder.

Usage:
    python generate_report.py --input data/hankamer_data_AY2627.xlsx
    python generate_report.py --input data/raw_data.csv
    python generate_report.py  # uses data/data_template.xlsx by default

Requirements:
    pip install -r requirements.txt
"""

import argparse
import math
import os
import shutil
import sys
from copy import deepcopy
from datetime import datetime
from pathlib import Path

import openpyxl
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree
from hex_map import generate_hex_map_png
from geo_tiles import generate_geo_tiles_png
import io

# ── paths ────────────────────────────────────────────────────────
BASE       = Path(__file__).parent
TEMPLATE   = BASE / "template.pptx"
OUTPUT_DIR = BASE / "output"

# ── helpers ──────────────────────────────────────────────────────

def load_data(path: Path) -> dict:
    """
    Load all sheets from the data file into a dict of DataFrames.
    Supports .xlsx and .csv.
    For CSV: treats the whole file as the 'attendance' sheet.
    """
    path = Path(path)
    if not path.exists():
        sys.exit(f"ERROR: Input file not found: {path}")

    ext = path.suffix.lower()
    if ext in (".xlsx", ".xls"):
        sheets = pd.read_excel(path, sheet_name=None, header=None)
        print(f"✔ Loaded {len(sheets)} sheets from {path.name}")
        return sheets
    elif ext == ".csv":
        df = pd.read_csv(path, header=None)
        print(f"✔ Loaded CSV from {path.name}")
        return {"attendance": df}
    else:
        sys.exit(f"ERROR: Unsupported file type '{ext}'. Use .xlsx or .csv")


def get_sheet(sheets: dict, name: str) -> pd.DataFrame:
    if name not in sheets:
        sys.exit(f"ERROR: Sheet '{name}' not found in data file. "
                 f"Available sheets: {list(sheets.keys())}")
    return sheets[name]


def cell(df: pd.DataFrame, row: int, col: int, default=None):
    """Get a value from a DataFrame by 0-based row/col, with a default."""
    try:
        v = df.iloc[row, col]
        if pd.isna(v):
            return default
        return v
    except (IndexError, KeyError):
        return default


def str_val(df, row, col, default=""):
    v = cell(df, row, col, default)
    return str(v).strip() if v is not None else default


def int_val(df, row, col, default=0):
    v = cell(df, row, col, default)
    try:
        return int(float(v))
    except (TypeError, ValueError):
        return default


def find_shape(slide, name: str):
    """Find a shape by name on a slide."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def set_text_run(shape, para_idx: int, new_text: str):
    """Replace text in a specific paragraph of a shape's text frame."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if para_idx >= len(tf.paragraphs):
        return
    para = tf.paragraphs[para_idx]
    if para.runs:
        # Preserve formatting of first run, replace text
        run = para.runs[0]
        run.text = new_text
        # Clear extra runs
        for extra in para.runs[1:]:
            extra.text = ""
    else:
        # Add a run if none exist
        from pptx.oxml.ns import qn
        r_elem = deepcopy(para._p)
        para._p.text = new_text


def replace_shape_text(slide, shape_name: str, new_text: str):
    """Replace ALL text in a shape with new_text, preserving paragraph/run formatting."""
    shape = find_shape(slide, shape_name)
    if not shape or not shape.has_text_frame:
        print(f"  ⚠ Shape '{shape_name}' not found or has no text")
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    if tf.paragraphs:
        if tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = new_text


def replace_para_text(slide, shape_name: str, para_idx: int, new_text: str):
    """Replace text in a specific paragraph index within a shape."""
    shape = find_shape(slide, shape_name)
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if para_idx >= len(tf.paragraphs):
        return
    para = tf.paragraphs[para_idx]
    if para.runs:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""


def update_table_cell(table, row: int, col: int, text: str):
    """Update a single cell in a python-pptx Table object."""
    cell_obj = table.cell(row, col)
    tf = cell_obj.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
    elif tf.paragraphs:
        from pptx.oxml.ns import qn
        p = tf.paragraphs[0]._p
        r = etree.SubElement(p, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'), attrib={'lang': 'en-US'})
        t = etree.SubElement(r, qn('a:t'))
        t.text = text


def update_chart_data(chart_shape, sheet_name: str, data_rows: list,
                       series_configs: list):
    """
    Directly update the embedded xlsx inside a chart shape.
    """
    from io import BytesIO

    # Find the embedded xlsx via the chart part's relationships
    chart_part = chart_shape.chart._part
    xlsx_part  = None
    for rel in chart_part.rels.values():
        if hasattr(rel._target, '_blob'):
            xlsx_part = rel._target
            break

    if xlsx_part is None:
        print(f"  ⚠ Could not find embedded workbook for chart '{chart_shape.name}'")
        return

    wb = openpyxl.load_workbook(BytesIO(xlsx_part._blob))

    # Get or rename the active sheet to match what the chart XML expects
    if sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
    else:
        ws = wb.active
        ws.title = sheet_name

    # Clear existing data area (rows 2 onward)
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row):
        for c in row:
            c.value = None

    # Write series labels in row 2
    for cfg in series_configs:
        ws.cell(row=cfg['label_row'], column=cfg['label_col'], value=cfg['label'])

    # Write data rows starting at row 3
    for i, row_data in enumerate(data_rows):
        r = i + 3
        ws.cell(row=r, column=1, value=row_data[0])   # category label
        for cfg in series_configs:
            val = row_data[cfg['col_idx']]
            ws.cell(row=r, column=cfg['data_col'],
                    value=val if val is not None else None)

    buf = BytesIO()
    wb.save(buf)
    xlsx_part._blob = buf.getvalue()
    print(f"  ✔ Chart '{chart_shape.name}' updated")


# ── hex map helpers ───────────────────────────────────────────────

def replace_hex_map(slide, state_counts, shape_name="Image 0", _png_override=None):
    """Replace the static image on slide 3 with an auto-generated one."""
    target = None
    for shape in slide.shapes:
        if shape.name == shape_name:
            target = shape
            break
    if target is None:
        raise ValueError(f"Shape '{shape_name}' not found. "
                         f"Available: {[s.name for s in slide.shapes]}")

    png_bytes = _png_override if _png_override else generate_hex_map_png(state_counts)
    left, top, width, height = target.left, target.top, target.width, target.height

    sp_tree = slide.shapes._spTree
    sp_tree.remove(target._element)

    pic = slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width, height)
    pic.name = shape_name
    return pic


# ── main ─────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Generate Hankamer info session report")
    parser.add_argument("--input", default="data/data_template.xlsx",
                        help="Path to input data file (.xlsx or .csv)")
    parser.add_argument("--output", default=None,
                        help="Output PPTX path (default: output/Hankamer_Report_YYYY-MM-DD.pptx)")
    args = parser.parse_args()

    input_path = Path(args.input)
    sheets = load_data(input_path)

    # ── read all sheets ──────────────────────────────────────────
    cfg_df  = get_sheet(sheets, "config")
    att_df  = get_sheet(sheets, "attendance")
    gy_df   = get_sheet(sheets, "grad_year")
    st_df   = get_sheet(sheets, "top10_states")
    all_st_df = sheets.get("all_states")  # optional — falls back to top10
    reg_df  = get_sheet(sheets, "regional")
    comp_df = get_sheet(sheets, "companion")
    bul_df  = get_sheet(sheets, "bullets")

    # ── parse config (col A = key, col B = value, starting row 3) ─
    config = {}
    for i in range(len(cfg_df)):
        k = str_val(cfg_df, i, 0)
        v = cell(cfg_df, i, 1)
        if k and not k.startswith("Field"):
            config[k] = v if v is not None else ""

    ay          = str(config.get("academic_year", "2025–2026"))
    fall_label  = str(config.get("fall_label",  "Fall 2025"))
    spring_label= str(config.get("spring_label","Spring 2026"))
    total_sess  = str(config.get("total_sessions", "117"))
    est_visitors= str(config.get("est_total_visitors", "~1,600"))
    states_repr = str(config.get("states_represented", "38"))
    top2_pct    = str(config.get("top2_classes_pct", "92%"))
    top2_label  = str(config.get("top2_classes_label", "Class of 2026 & 2027"))
    vis_mult    = str(config.get("visitor_multiplier", "2.5×"))
    footnote4   = str(config.get("footnote_slide4", ""))
    footer_text = str(config.get("footer_text", ""))

    # ── parse attendance ─────────────────────────────────────────
    # Rows 3–11 (0-based: 2–10), cols: A=month, B=fall, C=spring
    att_rows = []
    fall_total   = 0
    spring_total = 0
    for i in range(2, 11):
        month  = str_val(att_df, i, 0)
        fv     = cell(att_df, i, 1)
        sv     = cell(att_df, i, 2)
        fv_int = None if (fv is None or str(fv).strip() == "") else int(float(fv))
        sv_int = None if (sv is None or str(sv).strip() == "") else int(float(sv))
        if fv_int: fall_total   += fv_int
        if sv_int: spring_total += sv_int
        att_rows.append((month, fv_int, sv_int))

    grand_total = fall_total + spring_total

    # Override config if attendance data present
    if fall_total:   config["fall_students"]   = fall_total
    if spring_total: config["spring_students"] = spring_total
    fall_students   = int(config.get("fall_students",   fall_total))
    spring_students = int(config.get("spring_students", spring_total))

    print(f"\nData summary:")
    print(f"  AY: {ay}  |  Fall: {fall_students}  |  Spring: {spring_students}  |  Total: {grand_total}")

    # ── parse grad year ──────────────────────────────────────────
    # Rows 3–7 (0-based: 2–6), cols: A=class, B=count
    gy_rows = []
    gy_total = 0
    for i in range(2, 7):
        cls = str_val(gy_df, i, 0)
        cnt = int_val(gy_df, i, 1, 0)
        gy_total += cnt
        gy_rows.append((cls, cnt))

    # Cross-check
    if grand_total != gy_total and grand_total > 0 and gy_total > 0:
        print(f"  ⚠ WARNING: Attendance total ({grand_total}) ≠ Grad Year total ({gy_total}). "
              f"Check your data.")

    # ── parse top10 states ───────────────────────────────────────
    # Rows 3–12 (0-based: 2–11), cols: A=abbr, B=count, C=name
    # Input is largest→smallest; chart needs smallest→largest (ascending)
    st_rows = []
    for i in range(2, 12):
        abbr  = str_val(st_df, i, 0)
        cnt   = int_val(st_df, i, 1, 0)
        name  = str_val(st_df, i, 2)
        if abbr:
            st_rows.append((abbr, cnt, name))
    # Sort ascending for chart (chart reads bottom→top)
    st_rows_asc = sorted(st_rows, key=lambda x: x[1])

    # ── parse all_states (for tile grid on slide 3)
    all_state_rows = []
    if all_st_df is not None:
        for i in range(2, 52):
            abbr = str_val(all_st_df, i, 0)
            cnt  = int_val(all_st_df, i, 1, 0)
            if abbr and cnt:
                all_state_rows.append((abbr, cnt))
    if not all_state_rows:
        all_state_rows = [(a, c) for a, c, _ in st_rows]  # fallback to top10

    # ── parse regional ───────────────────────────────────────────
    # Rows 3–8 (0-based: 2–7): A=region, B=count
    reg_rows = []
    reg_total = 0
    for i in range(2, 8):
        reg   = str_val(reg_df, i, 0)
        cnt   = int_val(reg_df, i, 1, 0)
        reg_total += cnt
        reg_rows.append((reg, cnt))

    tx_cnt    = reg_rows[0][1] if reg_rows else 0
    intl_cnt  = reg_rows[5][1] if len(reg_rows) > 5 else 0
    oos_cnt   = reg_total - tx_cnt - intl_cnt

    tx_pct_str  = f"{round(tx_cnt  / reg_total * 100)}%"
    oos_pct_str = f"{round(oos_cnt / reg_total * 100)}%"
    intl_str    = str(intl_cnt)

    # ── parse companion ──────────────────────────────────────────
    # Rows 3–9 (0-based: 2–8): A=type, B=frequency
    comp_rows = []
    for i in range(2, 9):
        ctype = str_val(comp_df, i, 0)
        freq  = str_val(comp_df, i, 1)
        if ctype:
            comp_rows.append((ctype, freq))

    # ── parse bullets ────────────────────────────────────────────
    # Rows 3 onward (0-based: 2+), skipping section headers (col A only)
    # Col A = key, Col B = text
    bullets = {}
    for i in range(len(bul_df)):
        k = str_val(bul_df, i, 0)
        v = str_val(bul_df, i, 1)
        if k and v and not k.isupper():  # skip section header rows
            bullets[k] = v

    # ── open template and modify ──────────────────────────────────
    print("\nBuilding presentation...")
    prs = Presentation(str(TEMPLATE))

    slide1 = prs.slides[0]
    slide2 = prs.slides[1]
    slide3 = prs.slides[2]
    slide4 = prs.slides[3]

    # ════════════════════════════════════════════════════════════
    # SLIDE 1
    # ════════════════════════════════════════════════════════════
    print("  Slide 1...")

    # Academic year
    replace_shape_text(slide1, "Text 4", f"Academic Year {ay}")

    # Stat boxes — fall students
    replace_shape_text(slide1, "Text 7",  str(fall_students))
    sh = find_shape(slide1, "Text 8")
    if sh and sh.has_text_frame:
        paras = sh.text_frame.paragraphs
        if len(paras) >= 1 and paras[0].runs: paras[0].runs[0].text = fall_label
        if len(paras) >= 2 and paras[1].runs: paras[1].runs[0].text = "Students"

    # Spring students
    replace_shape_text(slide1, "Text 10", str(spring_students))
    sh = find_shape(slide1, "Text 11")
    if sh and sh.has_text_frame:
        paras = sh.text_frame.paragraphs
        if len(paras) >= 1 and paras[0].runs: paras[0].runs[0].text = spring_label
        if len(paras) >= 2 and paras[1].runs: paras[1].runs[0].text = "Students"

    # Sessions, visitors, states
    replace_shape_text(slide1, "Text 13", str(total_sess))
    replace_shape_text(slide1, "Text 16", str(est_visitors))
    replace_shape_text(slide1, "Text 19", str(states_repr))

    # Overview bullets (Text 23)
    ov_keys = ["overview_1","overview_2","overview_3","overview_4","overview_5","overview_6"]
    sh = find_shape(slide1, "Text 23")
    if sh and sh.has_text_frame:
        paras = sh.text_frame.paragraphs
        for i, key in enumerate(ov_keys):
            if i < len(paras) and key in bullets:
                if paras[i].runs:
                    paras[i].runs[0].text = bullets[key]

    # Recommendations bullets (Text 25)
    rec_keys = ["rec_1","rec_2","rec_3"]
    sh = find_shape(slide1, "Text 25")
    if sh and sh.has_text_frame:
        paras = sh.text_frame.paragraphs
        for i, key in enumerate(rec_keys):
            if i < len(paras) and key in bullets:
                if paras[i].runs:
                    paras[i].runs[0].text = bullets[key]

    # ════════════════════════════════════════════════════════════
    # SLIDE 2
    # ════════════════════════════════════════════════════════════
    print("  Slide 2...")

    # AY header
    replace_shape_text(slide2, "Text 2", f"AY {ay}")

    # Footnote
    if "slide2_footnote" in bullets:
        replace_shape_text(slide2, "Text 4", bullets["slide2_footnote"])

    # Callout box: "Class of 2026 & 2027 = 92% of students"
    callout_text = f"{top2_label} = {top2_pct} of students"
    replace_shape_text(slide2, "Text 8", callout_text)

    # ── Bar chart (Chart 0 on slide 2) ──────────────────────────
    bar_chart_shape = find_shape(slide2, "Chart 0")
    if bar_chart_shape:
        update_chart_data(
            bar_chart_shape,
            sheet_name="Attendance_Monthly",
            data_rows=att_rows,
            series_configs=[
                {"col_idx": 1, "data_col": 2, "label": fall_label,
                 "label_row": 2, "label_col": 2},
                {"col_idx": 2, "data_col": 3, "label": spring_label,
                 "label_row": 2, "label_col": 3},
            ]
        )

    # ── Pie chart (Chart 2 on slide 2) ──────────────────────────
    pie_chart_shape = find_shape(slide2, "Chart 2")
    if pie_chart_shape:
        pie_rows = [(cls, cnt, None) for cls, cnt in gy_rows]
        update_chart_data(
            pie_chart_shape,
            sheet_name="Grad_Year",
            data_rows=pie_rows,
            series_configs=[
                {"col_idx": 1, "data_col": 2, "label": "Grad Year",
                 "label_row": 2, "label_col": 2},
            ]
        )

    # ════════════════════════════════════════════════════════════
    # SLIDE 3
    # ════════════════════════════════════════════════════════════
    print("  Slide 3...")

    # Header — "GEOGRAPHIC REACH — AY 2025–2026"
    replace_shape_text(slide3, "Text 1", f"GEOGRAPHIC REACH — AY {ay}")

    # States represented badge
    replace_shape_text(slide3, "Text 2", f"{states_repr} States Represented")

    # Callout boxes
    replace_shape_text(slide3, "Text 7",  tx_pct_str)
    replace_shape_text(slide3, "Text 10", oos_pct_str)
    replace_shape_text(slide3, "Text 13", intl_str)

    # Regional table (Table 0)
    tbl_shape = find_shape(slide3, "Table 0")
    if tbl_shape:
        tbl = tbl_shape.table
        for i, (reg, cnt) in enumerate(reg_rows):
            row_idx = i + 1  # skip header row 0
            if row_idx < len(tbl.rows):
                pct = cnt / reg_total if reg_total else 0
                pct_str = "<1%" if 0 < pct < 0.005 else f"{round(pct * 100)}%"
                update_table_cell(tbl, row_idx, 0, reg)
                update_table_cell(tbl, row_idx, 1, str(cnt))
                update_table_cell(tbl, row_idx, 2, pct_str)

    # Geo bullets (Text 15)
    sh = find_shape(slide3, "Text 15")
    if sh and sh.has_text_frame:
        paras = sh.text_frame.paragraphs
        for i, key in enumerate(["geo_bullet_1", "geo_bullet_2"]):
            if i < len(paras) and key in bullets:
                if paras[i].runs:
                    paras[i].runs[0].text = bullets[key]

    # Geo footnote (Text 16)
    if "slide3_footnote" in bullets:
        replace_shape_text(slide3, "Text 16", bullets["slide3_footnote"])

    # ════════════════════════════════════════════════════════════
    # SLIDE 4
    # ════════════════════════════════════════════════════════════
    print("  Slide 4...")

    # Companion table (Table 0)
    tbl_shape = find_shape(slide4, "Table 0")
    if tbl_shape:
        tbl = tbl_shape.table
        for i, (ctype, freq) in enumerate(comp_rows):
            row_idx = i + 1  # skip header
            if row_idx < len(tbl.rows):
                update_table_cell(tbl, row_idx, 0, ctype)
                update_table_cell(tbl, row_idx, 1, freq)

    # Family engagement bullets (Text 9)
    sh = find_shape(slide4, "Text 9")
    if sh and sh.has_text_frame:
        paras = sh.text_frame.paragraphs
        for i, key in enumerate(["engage_1","engage_2","engage_3"]):
            if i < len(paras) and key in bullets:
                if paras[i].runs:
                    paras[i].runs[0].text = bullets[key]

    # Slide 4 footnote (Text 10)
    if footnote4:
        replace_shape_text(slide4, "Text 10", footnote4)

    # Footer (Text 12)
    if footer_text:
        replace_shape_text(slide4, "Text 12", footer_text)

    # ── Top 10 States chart (Chart 0 on slide 4) ────────────────
    states_chart_shape = find_shape(slide4, "Chart 0")
    if states_chart_shape:
        st_chart_rows = [(abbr, cnt, None) for abbr, cnt, _ in st_rows_asc]
        update_chart_data(
            states_chart_shape,
            sheet_name="Top10_States",
            data_rows=st_chart_rows,
            series_configs=[
                {"col_idx": 1, "data_col": 2, "label": "Students",
                 "label_row": 2, "label_col": 2},
            ]
        )

    # ── save output ──────────────────────────────────────────────
    OUTPUT_DIR.mkdir(exist_ok=True)
    if args.output:
        out_path = Path(args.output)
    else:
        safe_ay  = ay.replace("–","-").replace("/","-")
        out_path = OUTPUT_DIR / f"Hankamer_Report_AY{safe_ay}.pptx"

    # Generate and insert geo tile grid on slide 3
    geo_png = generate_geo_tiles_png(all_state_rows)
    replace_hex_map(prs.slides[2], {}, _png_override=geo_png)

    prs.save(str(out_path))
    print(f"\n✔ Report saved to: {out_path}")
    print(f"  Slides: 4  |  AY: {ay}  |  Students: {grand_total}")


if __name__ == "__main__":
    main()
