"""
app.py  —  Hankamer Report Generator Web App
Run locally:  streamlit run app.py
"""

import io
import sys
import traceback
from pathlib import Path

import pandas as pd
import streamlit as st

# ── page config ──────────────────────────────────────────────────
st.set_page_config(
    page_title="Hankamer Report Generator",
    page_icon="📊",
    layout="centered",
)

# ── branding ─────────────────────────────────────────────────────
st.markdown("""
<style>
    .main { background-color: #f9f9f9; }
    .stButton > button {
        background-color: #1A4731;
        color: white;
        font-size: 16px;
        padding: 12px 28px;
        border-radius: 6px;
        border: none;
        width: 100%;
    }
    .stButton > button:hover {
        background-color: #F5B800;
        color: #1A4731;
    }
    .stDownloadButton > button {
        background-color: #F5B800;
        color: #1A4731;
        font-size: 16px;
        font-weight: bold;
        padding: 12px 28px;
        border-radius: 6px;
        border: none;
        width: 100%;
    }
    h1 { color: #1A4731; }
    h3 { color: #1A4731; }
</style>
""", unsafe_allow_html=True)

st.image("https://www.baylor.edu/marketing/doc.php/233034.png", width=180) if False else None

col1, col2 = st.columns([1, 4])
with col1:
    st.markdown("## 🎓")
with col2:
    st.markdown("## Hankamer Report Generator")
    st.markdown("*Hankamer School of Business · Baylor University*")

st.divider()

# ── instructions ─────────────────────────────────────────────────
with st.expander("📋  How to use this tool", expanded=False):
    st.markdown("""
**Step 1 —** Fill in your `data_template.xlsx` file with this year's data (yellow cells only).

**Step 2 —** Upload it using the button below.

**Step 3 —** Click **Generate Report**.

**Step 4 —** Download the finished PowerPoint.

---

**Sheet names required in your Excel file:**

| Sheet | Contents |
|-------|----------|
| `config` | Academic year, session count, visitor estimate, etc. |
| `attendance` | Monthly headcount by semester |
| `grad_year` | Student count by graduation class |
| `top10_states` | Top 10 states (largest first) |
| `regional` | Regional breakdown |
| `companion` | Who joins students (frequency labels) |
| `bullets` | All editable text / bullet points |
""")

# ── file upload ───────────────────────────────────────────────────
st.markdown("### 1.  Upload your data file")
uploaded_file = st.file_uploader(
    "Accepts .xlsx or .csv",
    type=["xlsx", "csv"],
    label_visibility="collapsed",
)

if uploaded_file:
    st.success(f"✔  Uploaded: **{uploaded_file.name}**")

    # Preview sheet names
    if uploaded_file.name.endswith(".xlsx"):
        try:
            xl = pd.ExcelFile(uploaded_file)
            sheets = xl.sheet_names
            uploaded_file.seek(0)
            required = {"config","attendance","grad_year",
                        "top10_states","all_states","regional","companion","bullets"}
            found    = set(sheets)
            missing  = required - found
            if missing:
                st.error(f"⚠  Missing sheet(s): **{', '.join(sorted(missing))}**  "
                         f"— check your tab names match exactly.")
                st.stop()
            else:
                st.markdown(f"✔  All required sheets found: "
                            f"`{'`, `'.join(sorted(found & required))}`")
        except Exception as e:
            st.warning(f"Could not preview sheets: {e}")
            uploaded_file.seek(0)

    # ── generate button ───────────────────────────────────────────
    st.markdown("### 2.  Generate the report")

    if st.button("🚀  Generate Report"):
        uploaded_file.seek(0)
        file_bytes = uploaded_file.read()

        progress = st.progress(0, text="Reading data...")
        log_box  = st.empty()
        logs     = []

        def log(msg):
            logs.append(msg)
            log_box.markdown("\n".join(f"- {l}" for l in logs))

        try:
            # ── import generation logic ───────────────────────────
            # We inline the core logic here so the app is self-contained.
            import openpyxl
            import openpyxl.utils
            from pptx import Presentation
            from pptx.util import Pt
            from lxml import etree
            from copy import deepcopy

            # ── load data ─────────────────────────────────────────
            progress.progress(10, text="Loading sheets...")
            buf = io.BytesIO(file_bytes)
            if uploaded_file.name.endswith(".csv"):
                sheets = {"attendance": pd.read_csv(buf, header=None)}
            else:
                sheets = pd.read_excel(buf, sheet_name=None, header=None)

            log(f"✔ Loaded {len(sheets)} sheets")

            # ── helpers ───────────────────────────────────────────
            def sv(df, r, c, default=""):
                try:
                    v = df.iloc[r, c]
                    return default if pd.isna(v) else str(v).strip()
                except:
                    return default

            def iv(df, r, c, default=0):
                try:
                    v = df.iloc[r, c]
                    return default if pd.isna(v) else int(float(v))
                except:
                    return default

            def cv(df, r, c, default=None):
                try:
                    v = df.iloc[r, c]
                    return default if pd.isna(v) else v
                except:
                    return default

            def find_shape(slide, name):
                for s in slide.shapes:
                    if s.name == name:
                        return s
                return None

            def set_text(slide, shape_name, text):
                s = find_shape(slide, shape_name)
                if not s or not s.has_text_frame:
                    return
                tf = s.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ""
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = text

            def set_para(slide, shape_name, para_idx, text):
                s = find_shape(slide, shape_name)
                if not s or not s.has_text_frame:
                    return
                paras = s.text_frame.paragraphs
                if para_idx < len(paras) and paras[para_idx].runs:
                    paras[para_idx].runs[0].text = text
                    for run in paras[para_idx].runs[1:]:
                        run.text = ""

            def set_table_cell(table, row, col, text):
                cell = table.cell(row, col)
                tf = cell.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = text
                elif tf.paragraphs:
                    from pptx.oxml.ns import qn
                    p = tf.paragraphs[0]._p
                    r_el = etree.SubElement(p, qn('a:r'))
                    etree.SubElement(r_el, qn('a:rPr'), attrib={'lang':'en-US'})
                    t_el = etree.SubElement(r_el, qn('a:t'))
                    t_el.text = text

            def update_chart(chart_shape, sheet_name, data_rows, series_configs):
                chart_part = chart_shape.chart._part
                xlsx_part  = None
                for rel in chart_part.rels.values():
                    if hasattr(rel._target, '_blob'):
                        xlsx_part = rel._target
                        break
                if not xlsx_part:
                    return
                wb = openpyxl.load_workbook(io.BytesIO(xlsx_part._blob))
                ws = wb[sheet_name] if sheet_name in wb.sheetnames else wb.active
                ws.title = sheet_name
                for row in ws.iter_rows(min_row=2, max_row=ws.max_row):
                    for c in row:
                        c.value = None
                for cfg in series_configs:
                    ws.cell(row=cfg['label_row'], column=cfg['label_col'],
                            value=cfg['label'])
                for i, row_data in enumerate(data_rows):
                    r = i + 3
                    ws.cell(row=r, column=1, value=row_data[0])
                    for cfg in series_configs:
                        val = row_data[cfg['col_idx']]
                        ws.cell(row=r, column=cfg['data_col'],
                                value=val if val is not None else None)
                out = io.BytesIO()
                wb.save(out)
                xlsx_part._blob = out.getvalue()

            # ── parse config ──────────────────────────────────────
            progress.progress(20, text="Parsing config...")
            cfg_df = sheets["config"]
            config = {}
            for i in range(len(cfg_df)):
                k = sv(cfg_df, i, 0)
                v = cv(cfg_df, i, 1)
                if k and not k.startswith("Field"):
                    config[k] = v if v is not None else ""

            ay           = str(config.get("academic_year",      "2025–2026"))
            fall_label   = str(config.get("fall_label",         "Fall 2025"))
            spring_label = str(config.get("spring_label",       "Spring 2026"))
            total_sess   = str(config.get("total_sessions",     "117"))
            est_visitors = str(config.get("est_total_visitors", "~1,600"))
            states_repr  = str(config.get("states_represented", "38"))
            top2_pct     = str(config.get("top2_classes_pct",   "92%"))
            top2_label   = str(config.get("top2_classes_label", "Class of 2026 & 2027"))
            footnote4    = str(config.get("footnote_slide4",    ""))
            footer_text  = str(config.get("footer_text",        ""))

            # ── parse attendance ──────────────────────────────────
            progress.progress(30, text="Parsing attendance...")
            att_df = sheets["attendance"]
            att_rows     = []
            fall_total   = 0
            spring_total = 0
            for i in range(2, 11):
                month = sv(att_df, i, 0)
                fv    = cv(att_df, i, 1)
                sv_   = cv(att_df, i, 2)
                fv_i  = None if (fv is None or str(fv).strip()=="") else int(float(fv))
                sv_i  = None if (sv_ is None or str(sv_).strip()=="") else int(float(sv_))
                if fv_i: fall_total   += fv_i
                if sv_i: spring_total += sv_i
                att_rows.append((month, fv_i, sv_i))

            grand_total    = fall_total + spring_total
            fall_students  = int(config.get("fall_students",   fall_total   or 0))
            spring_students= int(config.get("spring_students", spring_total or 0))
            if fall_total:   fall_students   = fall_total
            if spring_total: spring_students = spring_total

            log(f"✔ Attendance: Fall {fall_students}, Spring {spring_students}, Total {grand_total}")

            # ── parse grad year ───────────────────────────────────
            gy_df  = sheets["grad_year"]
            gy_rows= []
            gy_total = 0
            for i in range(2, 7):
                cls = sv(gy_df, i, 0)
                cnt = iv(gy_df, i, 1, 0)
                gy_total += cnt
                gy_rows.append((cls, cnt))

            if grand_total != gy_total and grand_total > 0 and gy_total > 0:
                st.warning(f"⚠ Attendance total ({grand_total}) ≠ Grad Year total "
                           f"({gy_total}). Check your data — continuing anyway.")

            # ── parse top10 states ────────────────────────────────
            st_df   = sheets["top10_states"]
            st_rows = []
            for i in range(2, 12):
                abbr = sv(st_df, i, 0)
                cnt  = iv(st_df, i, 1, 0)
                if abbr:
                    st_rows.append((abbr, cnt))
            st_rows_asc = sorted(st_rows, key=lambda x: x[1])

            # ── parse regional ────────────────────────────────────
            reg_df   = sheets["regional"]
            reg_rows = []
            reg_total = 0
            for i in range(2, 8):
                reg = sv(reg_df, i, 0)
                cnt = iv(reg_df, i, 1, 0)
                reg_total += cnt
                reg_rows.append((reg, cnt))

            tx_cnt   = reg_rows[0][1] if reg_rows else 0
            intl_cnt = reg_rows[5][1] if len(reg_rows) > 5 else 0
            oos_cnt  = reg_total - tx_cnt - intl_cnt
            tx_pct_str  = f"{round(tx_cnt  / reg_total * 100)}%" if reg_total else "0%"
            oos_pct_str = f"{round(oos_cnt / reg_total * 100)}%" if reg_total else "0%"
            intl_str    = str(intl_cnt)

            # ── parse all_states ─────────────────────────────
            all_states_df = sheets.get("all_states")
            all_state_rows = []
            if all_states_df is not None:
                for i in range(2, len(all_states_df)):
                    abbr = sv(all_states_df, i, 0)
                    cnt  = iv(all_states_df, i, 1, 0)
                    if abbr and cnt:
                        all_state_rows.append((abbr, cnt))
            if not all_state_rows:
                all_state_rows = st_rows  # fallback to top10

            # ── parse companion ───────────────────────────────────
            comp_df   = sheets["companion"]
            comp_rows = []
            for i in range(2, 9):
                ctype = sv(comp_df, i, 0)
                freq  = sv(comp_df, i, 1)
                if ctype:
                    comp_rows.append((ctype, freq))

            # ── parse bullets ─────────────────────────────────────
            bul_df  = sheets["bullets"]
            bullets = {}
            for i in range(len(bul_df)):
                k = sv(bul_df, i, 0)
                v = sv(bul_df, i, 1)
                if k and v and not k.isupper():
                    bullets[k] = v

            log(f"✔ Data parsed — {len(bullets)} text blocks loaded")

            # ── open template ─────────────────────────────────────
            progress.progress(50, text="Opening template...")
            template_path = Path(__file__).parent / "template.pptx"
            if not template_path.exists():
                st.error("❌ template.pptx not found in the app folder.")
                st.stop()

            prs   = Presentation(str(template_path))
            slide1 = prs.slides[0]
            slide2 = prs.slides[1]
            slide3 = prs.slides[2]
            slide4 = prs.slides[3]

            # ── SLIDE 1 ───────────────────────────────────────────
            progress.progress(60, text="Building slide 1...")
            set_text(slide1, "Text 4", f"Academic Year {ay}")
            set_text(slide1, "Text 7",  str(fall_students))
            sh = find_shape(slide1, "Text 8")
            if sh and sh.has_text_frame:
                paras = sh.text_frame.paragraphs
                if len(paras) >= 1 and paras[0].runs: paras[0].runs[0].text = fall_label
                if len(paras) >= 2 and paras[1].runs: paras[1].runs[0].text = "Students"
            set_text(slide1, "Text 10", str(spring_students))
            sh = find_shape(slide1, "Text 11")
            if sh and sh.has_text_frame:
                paras = sh.text_frame.paragraphs
                if len(paras) >= 1 and paras[0].runs: paras[0].runs[0].text = spring_label
                if len(paras) >= 2 and paras[1].runs: paras[1].runs[0].text = "Students"
            set_text(slide1, "Text 13", str(total_sess))
            set_text(slide1, "Text 16", str(est_visitors))
            set_text(slide1, "Text 19", str(states_repr))

            ov_keys  = ["overview_1","overview_2","overview_3",
                        "overview_4","overview_5","overview_6"]
            sh = find_shape(slide1, "Text 23")
            if sh and sh.has_text_frame:
                for i, key in enumerate(ov_keys):
                    if i < len(sh.text_frame.paragraphs) and key in bullets:
                        p = sh.text_frame.paragraphs[i]
                        if p.runs: p.runs[0].text = bullets[key]

            rec_keys = ["rec_1","rec_2","rec_3"]
            sh = find_shape(slide1, "Text 25")
            if sh and sh.has_text_frame:
                for i, key in enumerate(rec_keys):
                    if i < len(sh.text_frame.paragraphs) and key in bullets:
                        p = sh.text_frame.paragraphs[i]
                        if p.runs: p.runs[0].text = bullets[key]

            log("✔ Slide 1 done")

            # ── SLIDE 2 ───────────────────────────────────────────
            progress.progress(70, text="Building slide 2 (charts)...")
            set_text(slide2, "Text 2", f"AY {ay}")
            if "slide2_footnote" in bullets:
                set_text(slide2, "Text 4", bullets["slide2_footnote"])
            set_text(slide2, "Text 8",
                     f"{top2_label} = {top2_pct} of students")

            bar = find_shape(slide2, "Chart 0")
            if bar:
                update_chart(bar, "Attendance_Monthly", att_rows, [
                    {"col_idx":1,"data_col":2,"label":fall_label,
                     "label_row":2,"label_col":2},
                    {"col_idx":2,"data_col":3,"label":spring_label,
                     "label_row":2,"label_col":3},
                ])

            pie = find_shape(slide2, "Chart 2")
            if pie:
                pie_rows = [(c, n, None) for c, n in gy_rows]
                update_chart(pie, "Grad_Year", pie_rows, [
                    {"col_idx":1,"data_col":2,"label":"Grad Year",
                     "label_row":2,"label_col":2},
                ])

            log("✔ Slide 2 done")

            # ── SLIDE 3 ───────────────────────────────────────────
            progress.progress(80, text="Building slide 3...")
            set_text(slide3, "Text 1",  f"GEOGRAPHIC REACH — AY {ay}")
            set_text(slide3, "Text 2",  f"{states_repr} States Represented")
            set_text(slide3, "Text 7",  tx_pct_str)
            set_text(slide3, "Text 10", oos_pct_str)
            set_text(slide3, "Text 13", intl_str)

            tbl_s = find_shape(slide3, "Table 0")
            if tbl_s:
                tbl = tbl_s.table
                for i, (reg, cnt) in enumerate(reg_rows):
                    ri = i + 1
                    if ri < len(tbl.rows):
                        pct = cnt / reg_total if reg_total else 0
                        pct_str = "<1%" if 0 < pct < 0.005 else f"{round(pct*100)}%"
                        set_table_cell(tbl, ri, 0, reg)
                        set_table_cell(tbl, ri, 1, str(cnt))
                        set_table_cell(tbl, ri, 2, pct_str)

            sh = find_shape(slide3, "Text 15")
            if sh and sh.has_text_frame:
                for i, key in enumerate(["geo_bullet_1","geo_bullet_2"]):
                    if i < len(sh.text_frame.paragraphs) and key in bullets:
                        p = sh.text_frame.paragraphs[i]
                        if p.runs: p.runs[0].text = bullets[key]

            if "slide3_footnote" in bullets:
                set_text(slide3, "Text 16", bullets["slide3_footnote"])

            log("✔ Slide 3 done")

            # ── SLIDE 4 ───────────────────────────────────────────
            progress.progress(90, text="Building slide 4...")
            tbl_s = find_shape(slide4, "Table 0")
            if tbl_s:
                tbl = tbl_s.table
                for i, (ctype, freq) in enumerate(comp_rows):
                    if i + 1 < len(tbl.rows):
                        set_table_cell(tbl, i+1, 0, ctype)
                        set_table_cell(tbl, i+1, 1, freq)

            sh = find_shape(slide4, "Text 9")
            if sh and sh.has_text_frame:
                for i, key in enumerate(["engage_1","engage_2","engage_3"]):
                    if i < len(sh.text_frame.paragraphs) and key in bullets:
                        p = sh.text_frame.paragraphs[i]
                        if p.runs: p.runs[0].text = bullets[key]

            if footnote4:
                set_text(slide4, "Text 10", footnote4)
            if footer_text:
                set_text(slide4, "Text 12", footer_text)

            st_chart = find_shape(slide4, "Chart 0")
            if st_chart:
                st_chart_rows = [(abbr, cnt, None) for abbr, cnt in st_rows_asc]
                update_chart(st_chart, "Top10_States", st_chart_rows, [
                    {"col_idx":1,"data_col":2,"label":"Students",
                     "label_row":2,"label_col":2},
                ])

            log("✔ Slide 4 done")

            # ── save to bytes ─────────────────────────────────────
            progress.progress(95, text="Saving presentation...")
            out_buf = io.BytesIO()
            prs.save(out_buf)
            out_buf.seek(0)
            pptx_bytes = out_buf.read()

            progress.progress(100, text="Done!")
            log(f"✔ Report complete — AY {ay}, {grand_total} students, 4 slides")

            # ── download button ───────────────────────────────────
            safe_ay   = ay.replace("–","-").replace("/","-")
            file_name = f"Hankamer_Report_AY{safe_ay}.pptx"

            st.success("✅  Report ready!")
            st.markdown("### 3.  Download your report")
            st.download_button(
                label=f"⬇  Download {file_name}",
                data=pptx_bytes,
                file_name=file_name,
                mime="application/vnd.openxmlformats-officedocument"
                     ".presentationml.presentation",
            )

        except Exception as e:
            progress.empty()
            st.error(f"❌  Something went wrong:\n\n```\n{traceback.format_exc()}\n```")
            st.info("Copy the error above and share it to get help fixing it.")

else:
    # No file uploaded yet — show placeholder
    st.info("👆  Upload your filled-in Excel file to get started.")
    st.markdown("Don't have the template? Download it below.")
    template_xlsx = Path(__file__).parent / "data" / "data_template.xlsx"
    if template_xlsx.exists():
        with open(template_xlsx, "rb") as f:
            st.download_button(
                label="⬇  Download data_template.xlsx",
                data=f.read(),
                file_name="data_template.xlsx",
                mime="application/vnd.openxmlformats-officedocument"
                     ".spreadsheetml.sheet",
            )

# ── footer ────────────────────────────────────────────────────────
st.divider()
st.caption("Hankamer School of Business · Baylor University · "
           "Internal Report Tool · For help, contact your administrator.")
